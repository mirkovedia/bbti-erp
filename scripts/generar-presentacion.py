"""
Genera la presentación ejecutiva del estado del proyecto BBTI ERP.

    python scripts/generar-presentacion.py [salida.pptx]

Por defecto escribe `BBTI-ERP-Estado.pptx` en la raíz del repositorio.
Los datos provienen de docs/ENTREGA.md: al actualizar el estado del proyecto,
actualizar ambos.

Requiere: pip install python-pptx
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Identidad visual (la del propio ERP: petróleo + ámbar de BBTI) ──────────
BG       = RGBColor(0x0B, 0x19, 0x20)   # petróleo profundo
SURFACE  = RGBColor(0x12, 0x24, 0x2D)   # paneles
SURFACE2 = RGBColor(0x18, 0x31, 0x3C)   # cabeceras de tabla
LINE     = RGBColor(0x23, 0x43, 0x50)   # hairlines
TEXT     = RGBColor(0xE7, 0xEF, 0xF2)
MUTED    = RGBColor(0x8B, 0xA7, 0xB3)
AMBER    = RGBColor(0xEC, 0x9D, 0x2E)   # acento de marca
TEAL     = RGBColor(0x4A, 0x92, 0xAD)
OK       = RGBColor(0x45, 0xB9, 0x8A)
WAIT     = RGBColor(0x71, 0x8E, 0x9B)

SANS = "Segoe UI"
MONO = "Consolas"

W, H = 13.333, 7.5          # 16:9
MARGIN = 0.85


def nueva():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    return prs


def lamina(prs):
    """Lámina en blanco con el fondo de marca."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fondo = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fondo.fill.solid()
    fondo.fill.fore_color.rgb = BG
    fondo.line.fill.background()
    fondo.shadow.inherit = False
    return s


def texto(slide, x, y, w, h, contenido, size=18, color=TEXT, bold=False,
          font=SANS, align=PP_ALIGN.LEFT, space=0, interlinea=1.25):
    caja = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = caja.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = interlinea
    r = p.add_run()
    r.text = contenido
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    if space:
        # python-pptx no expone letter-spacing: se emula vía XML
        r.font._rPr.set("spc", str(int(space * 100)))
    return caja


def eyebrow(slide, contenido, y=0.62):
    texto(slide, MARGIN, y, 10, 0.3, contenido.upper(), size=11,
          color=AMBER, font=MONO, space=1.6)


def titulo(slide, contenido, y=1.02, size=36, w=11.2):
    """Título de lámina. Debe caber en UNA línea: si envuelve, pisa el subtítulo.
    El aviso evita que un cambio de texto rompa la maqueta sin que se note."""
    por_linea = int((w * 72) / (size * 0.52))   # ancho medio de glifo ≈ 0.52 em
    if len(contenido) > por_linea:
        print(f"  AVISO: título de {len(contenido)} caracteres supera los "
              f"~{por_linea} que caben en una línea → «{contenido}»")
    texto(slide, MARGIN, y, w, 1.1, contenido, size=size, bold=True, interlinea=1.05)


def panel(slide, x, y, w, h, relleno=SURFACE, borde=LINE):
    p = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    p.fill.solid()
    p.fill.fore_color.rgb = relleno
    p.line.color.rgb = borde
    p.line.width = Pt(0.75)
    p.shadow.inherit = False
    return p


def regla(slide, y, x=MARGIN, w=W - 2 * MARGIN, color=LINE):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Emu(9525))
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    return r


def pie(slide, numero):
    texto(slide, MARGIN, H - 0.62, 8, 0.3, "BBTI ERP · Estado del proyecto",
          size=9, color=WAIT, font=MONO, space=1.2)
    texto(slide, W - MARGIN - 1.2, H - 0.62, 1.2, 0.3, f"{numero:02d}",
          size=9, color=WAIT, font=MONO, align=PP_ALIGN.RIGHT)


# ── Láminas ────────────────────────────────────────────────────────────────

def portada(prs):
    s = lamina(prs)
    # banda de acento vertical
    banda = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.14), prs.slide_height)
    banda.fill.solid()
    banda.fill.fore_color.rgb = AMBER
    banda.line.fill.background()
    banda.shadow.inherit = False

    texto(s, MARGIN, 2.35, 10, 0.35, "BBTI · SISTEMA DE GESTIÓN DE PROYECTOS",
          size=12, color=AMBER, font=MONO, space=2.0)
    texto(s, MARGIN, 2.95, 11.2, 1.9,
          "El ERP está completo, verificado\ny listo para desplegar.",
          size=42, bold=True, interlinea=1.08)
    texto(s, MARGIN, 4.85, 8.6, 1.0,
          "Cubre el ciclo íntegro del proyecto — de la orden comercial al cierre "
          "con pago al 100% — y ya no depende de ningún servicio de terceros.",
          size=15, color=MUTED, interlinea=1.45)
    regla(s, 6.15, w=3.2, color=LINE)
    texto(s, MARGIN, 6.4, 6, 0.35, "3 de agosto de 2026", size=11,
          color=WAIT, font=MONO, space=1.2)
    return s


def resumen(prs, n):
    s = lamina(prs)
    eyebrow(s, "Resumen ejecutivo")
    titulo(s, "Dónde está el proyecto")

    texto(s, MARGIN, 2.15, 11.4, 0.9,
          "El sistema funciona hoy de forma completa y demostrable. "
          "Lo único que resta es contratar un servidor y un dominio.",
          size=17, color=TEXT, interlinea=1.4)

    datos = [
        ("190+", "verificaciones automatizadas,\ntodas en verde"),
        ("0", "servicios externos de los\nque depende"),
        ("21", "tablas de datos\nen producción"),
        ("~1 h", "para estar en línea con\nservidor y dominio"),
    ]
    ancho = (W - 2 * MARGIN - 0.45) / 4
    for i, (cifra, etiqueta) in enumerate(datos):
        x = MARGIN + i * (ancho + 0.15)
        panel(s, x, 3.5, ancho, 1.75)
        texto(s, x + 0.28, 3.78, ancho - 0.5, 0.6, cifra, size=32, bold=True,
              font=MONO, color=AMBER if i == 0 else TEXT)
        texto(s, x + 0.28, 4.42, ancho - 0.5, 0.7, etiqueta, size=11,
              color=MUTED, interlinea=1.3)

    panel(s, MARGIN, 5.62, W - 2 * MARGIN, 0.72, relleno=SURFACE2, borde=LINE)
    punto = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(MARGIN + 0.3), Inches(5.9),
                               Inches(0.14), Inches(0.14))
    punto.fill.solid()
    punto.fill.fore_color.rgb = OK
    punto.line.fill.background()
    punto.shadow.inherit = False
    texto(s, MARGIN + 0.58, 5.85, 10, 0.3, "Sistema operativo y demostrable hoy",
          size=13, color=OK, bold=True)
    pie(s, n)
    return s


def alcance(prs, n):
    s = lamina(prs)
    eyebrow(s, "Alcance")
    titulo(s, "Qué hace el sistema")
    texto(s, MARGIN, 1.98, 10.6, 0.5,
          "Cada área ve el avance de las demás, pero solo edita lo suyo. "
          "El estado del proyecto se deriva de las firmas de cada etapa.",
          size=13, color=MUTED, interlinea=1.4)

    areas = [
        ("Comercial", "Crea la orden, registra monto y fecha de entrega, sube OC y comprobantes, importa el metrado desde Excel."),
        ("Ingeniería", "Sube y versiona planos, controla su estado de aprobación y deja observaciones técnicas."),
        ("Logística", "Gestiona los materiales del metrado y el avance de compras hasta completarlas."),
        ("Producción", "Controla las 7 etapas de fabricación. Cada una registra quién la completó, con fecha y hora."),
        ("Finanzas", "Registra pagos y autoriza el cierre: sin el 100% cobrado, el proyecto no se completa."),
        ("Transversal", "Panel de control, calendario, notificaciones, bitácora, papelera y alertas de vencimiento."),
    ]
    ancho = (W - 2 * MARGIN - 0.3) / 3
    alto = 1.75
    for i, (nombre, desc) in enumerate(areas):
        col, fila = i % 3, i // 3
        x = MARGIN + col * (ancho + 0.15)
        y = 2.75 + fila * (alto + 0.22)
        panel(s, x, y, ancho, alto)
        marca = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.28), Inches(y + 0.38),
                                   Inches(0.22), Emu(22000))
        marca.fill.solid()
        marca.fill.fore_color.rgb = AMBER
        marca.line.fill.background()
        marca.shadow.inherit = False
        texto(s, x + 0.62, y + 0.26, ancho - 0.9, 0.35, nombre, size=15, bold=True)
        texto(s, x + 0.28, y + 0.72, ancho - 0.56, 0.9, desc, size=10.5,
              color=MUTED, interlinea=1.35)
    pie(s, n)
    return s


def migracion(prs, n):
    s = lamina(prs)
    eyebrow(s, "Migración · 18 tareas")
    titulo(s, "Del alquiler a la infraestructura propia")
    texto(s, MARGIN, 2.05, 10.8, 0.5,
          "Cada pieza alquilada se reemplazó por una equivalente bajo control propio, "
          "sin perder funcionalidad.", size=13, color=MUTED, interlinea=1.4)

    filas = [
        ("Base de datos", "Supabase", "PostgreSQL con Prisma"),
        ("Autenticación", "Supabase Auth", "Propia · sesión cifrada y bcrypt"),
        ("Archivos", "Supabase Storage", "Almacenamiento privado con enlaces firmados"),
        ("Tiempo real", "Supabase Realtime", "Actualización automática propia"),
        ("Hosting", "Vercel", "Docker en cualquier servidor Linux"),
        ("Tareas diarias", "Vercel Cron", "Contenedor propio en el despliegue"),
    ]
    x0, y0 = MARGIN, 2.78
    ancho_total = W - 2 * MARGIN
    c1, c2 = 2.9, 3.4
    c3 = ancho_total - c1 - c2
    alto_fila = 0.52

    # cabecera
    panel(s, x0, y0, ancho_total, 0.42, relleno=SURFACE2, borde=LINE)
    for etiqueta, dx, dw in (("COMPONENTE", 0.22, c1), ("ANTES", c1 + 0.22, c2), ("AHORA", c1 + c2 + 0.22, c3)):
        texto(s, x0 + dx, y0 + 0.12, dw, 0.25, etiqueta, size=9.5, color=MUTED,
              font=MONO, space=1.5)

    for i, (comp, antes, ahora) in enumerate(filas):
        y = y0 + 0.42 + i * alto_fila
        panel(s, x0, y, ancho_total, alto_fila, relleno=SURFACE, borde=LINE)
        texto(s, x0 + 0.22, y + 0.14, c1 - 0.3, 0.3, comp, size=12, bold=True)
        texto(s, x0 + c1 + 0.22, y + 0.15, c2 - 0.3, 0.3, antes, size=11.5, color=MUTED)
        texto(s, x0 + c1 + c2 + 0.22, y + 0.15, 0.3, 0.3, "→", size=12,
              color=AMBER, font=MONO, bold=True)
        texto(s, x0 + c1 + c2 + 0.55, y + 0.15, c3 - 0.7, 0.3, ahora, size=11.5)
    pie(s, n)
    return s


def lista_dos_columnas(prs, n, ojo, tit, subtitulo, items, size_item=11.5):
    s = lamina(prs)
    eyebrow(s, ojo)
    titulo(s, tit)
    y_inicio = 2.15
    if subtitulo:
        texto(s, MARGIN, 1.98, 10.8, 0.5, subtitulo, size=13, color=MUTED, interlinea=1.4)
        y_inicio = 2.6

    ancho = (W - 2 * MARGIN - 0.6) / 2
    mitad = (len(items) + 1) // 2
    for i, (fuerte, resto) in enumerate(items):
        col, fila = (0, i) if i < mitad else (1, i - mitad)
        x = MARGIN + col * (ancho + 0.6)
        y = y_inicio + fila * 0.78
        rombo = s.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(x), Inches(y + 0.12),
                                   Inches(0.11), Inches(0.11))
        rombo.fill.solid()
        rombo.fill.fore_color.rgb = AMBER
        rombo.line.fill.background()
        rombo.shadow.inherit = False

        caja = s.shapes.add_textbox(Inches(x + 0.3), Inches(y), Inches(ancho - 0.3), Inches(0.72))
        tf = caja.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.line_spacing = 1.3
        r1 = p.add_run()
        r1.text = fuerte + " "
        r1.font.size = Pt(size_item)
        r1.font.bold = True
        r1.font.name = SANS
        r1.font.color.rgb = TEXT
        r2 = p.add_run()
        r2.text = resto
        r2.font.size = Pt(size_item)
        r2.font.name = SANS
        r2.font.color.rgb = MUTED
    pie(s, n)
    return s


def verificacion(prs, n):
    s = lamina(prs)
    eyebrow(s, "Control de calidad")
    titulo(s, "Cómo sabemos que funciona")
    texto(s, MARGIN, 2.05, 10.8, 0.5,
          "Las pruebas no simulan el sistema: lo ejercitan de verdad, con los usuarios "
          "reales, sobre la base de datos y el almacenamiento reales.",
          size=13, color=MUTED, interlinea=1.4)

    filas = [
        ("Simulación multi-usuario", "Los 7 usuarios ejercitando el ciclo completo: permisos por rol, firmas, pagos y avisos", "30 / 30"),
        ("Seguridad", "Revocación de sesiones, contraseñas, bloqueos, cabeceras y bitácora", "27 / 27"),
        ("Flujo de negocio", "De la creación del proyecto al cierre con el pago al 100%", "8 / 8"),
        ("Firmas de etapa", "Confirmaciones, validaciones, permisos y reversión en cascada", "12 / 12"),
        ("Lógica de negocio", "Estados, vencimientos, cálculo de pagos y sesiones", "75 / 75"),
        ("Funcionalidades", "Documentos, metrado, notificaciones, productividad, alertas y actividad", "en verde"),
    ]
    x0, y0 = MARGIN, 2.85
    ancho_total = W - 2 * MARGIN
    c1, c3 = 3.0, 1.5
    c2 = ancho_total - c1 - c3
    alto_fila = 0.53

    panel(s, x0, y0, ancho_total, 0.42, relleno=SURFACE2, borde=LINE)
    for etiqueta, dx, dw, al in (("PRUEBA", 0.22, c1, PP_ALIGN.LEFT),
                                 ("QUÉ COMPRUEBA", c1 + 0.22, c2, PP_ALIGN.LEFT),
                                 ("RESULTADO", c1 + c2, c3 - 0.25, PP_ALIGN.RIGHT)):
        texto(s, x0 + dx, y0 + 0.12, dw, 0.25, etiqueta, size=9.5, color=MUTED,
              font=MONO, space=1.5, align=al)

    for i, (prueba, que, res) in enumerate(filas):
        y = y0 + 0.42 + i * alto_fila
        panel(s, x0, y, ancho_total, alto_fila, relleno=SURFACE, borde=LINE)
        texto(s, x0 + 0.22, y + 0.15, c1 - 0.3, 0.3, prueba, size=11.5, bold=True)
        texto(s, x0 + c1 + 0.22, y + 0.16, c2 - 0.3, 0.3, que, size=10.5, color=MUTED)
        texto(s, x0 + c1 + c2, y + 0.15, c3 - 0.25, 0.3, res, size=12,
              color=OK, font=MONO, bold=True, align=PP_ALIGN.RIGHT)
    pie(s, n)
    return s


def estado(prs, n):
    s = lamina(prs)
    eyebrow(s, "Situación")
    titulo(s, "Estado actual")

    bloques = [
        (OK, "Completado", "Sistema construido y verificado",
         "Aplicación completa y migrada · empaquetada en Docker ·\nseguridad comprobada · documentación de despliegue"),
        (AMBER, "En curso", "Autonomía total del despliegue",
         "La base de datos y los archivos pasarán a instalarse junto a la\naplicación, con respaldos diarios automáticos. Diseño aprobado."),
        (WAIT, "Pendiente", "Puesta en producción",
         "Contratar servidor y dominio. No requiere desarrollo:\nel sistema ya está listo para instalarse."),
    ]
    y = 2.2
    for color, etiqueta, tit, desc in bloques:
        nodo = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(MARGIN), Inches(y + 0.14),
                                  Inches(0.2), Inches(0.2))
        nodo.fill.solid()
        nodo.fill.fore_color.rgb = color
        nodo.line.fill.background()
        nodo.shadow.inherit = False
        # riel que conecta los nodos (secuencia real)
        if etiqueta != "Pendiente":
            riel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN + 0.09),
                                      Inches(y + 0.34), Emu(9525), Inches(1.32))
            riel.fill.solid()
            riel.fill.fore_color.rgb = LINE
            riel.line.fill.background()
            riel.shadow.inherit = False

        texto(s, MARGIN + 0.45, y + 0.05, 5.0, 0.4, tit, size=17, bold=True)
        # etiqueta junto al título, no al borde: mantiene la lectura en un solo bloque
        panel(s, MARGIN + 5.35, y + 0.08, 1.55, 0.32, relleno=BG, borde=color)
        texto(s, MARGIN + 5.35, y + 0.13, 1.55, 0.25, etiqueta.upper(), size=9,
              color=color, font=MONO, align=PP_ALIGN.CENTER, space=1.2)
        texto(s, MARGIN + 0.45, y + 0.52, 9.5, 0.75, desc, size=11.5,
              color=MUTED, interlinea=1.35)
        y += 1.58
    pie(s, n)
    return s


def pendiente(prs, n):
    s = lamina(prs)
    eyebrow(s, "Lo único que falta")
    titulo(s, "Puesta en producción")
    texto(s, MARGIN, 2.05, 10.8, 0.5,
          "No requiere desarrollo, solo contratación. Con el servidor y el dominio "
          "listos, el sistema queda en línea en aproximadamente una hora.",
          size=13, color=MUTED, interlinea=1.4)

    items = [
        ("Servidor", "≈ 6–10 USD / mes", "Cualquier proveedor\ncon Docker · 30 min"),
        ("Dominio", "≈ 10–15 USD / año", "Registro anual\n· 15 min"),
        ("Despliegue", "4 comandos", "Certificado HTTPS\nautomático · 30 min"),
    ]
    ancho = (W - 2 * MARGIN - 0.4) / 3
    for i, (tit, costo, det) in enumerate(items):
        x = MARGIN + i * (ancho + 0.2)
        panel(s, x, 2.95, ancho, 2.0)
        texto(s, x + 0.32, 3.22, ancho - 0.6, 0.4, tit, size=18, bold=True)
        texto(s, x + 0.32, 3.75, ancho - 0.6, 0.35, costo, size=14,
              color=AMBER, font=MONO, bold=True)
        texto(s, x + 0.32, 4.22, ancho - 0.6, 0.6, det, size=11,
              color=MUTED, interlinea=1.35)

    panel(s, MARGIN, 5.35, W - 2 * MARGIN, 1.05, relleno=SURFACE2, borde=LINE)
    texto(s, MARGIN + 0.35, 5.55, 11.5, 0.3, "Sin quedar atado a un proveedor",
          size=13, bold=True, color=AMBER)
    texto(s, MARGIN + 0.35, 5.86, 11.4, 0.35,
          "Todo se apoya en estándares abiertos. Mover la base de datos o los archivos "
          "a un servicio en la nube es cambiar configuración, no reescribir el sistema.",
          size=11, color=MUTED)
    pie(s, n)
    return s


def despliegue(prs, n):
    s = lamina(prs)
    eyebrow(s, "Puesta en marcha")
    titulo(s, "Del código al sistema en línea")

    ancho_izq = 7.0
    panel(s, MARGIN, 2.35, ancho_izq, 3.75)   # holgura bajo la última línea
    comandos = [
        ("# 1 · Instalar Docker en el servidor", MUTED),
        ("curl -fsSL https://get.docker.com | sh", TEXT),
        ("", TEXT),
        ("# 2 · Obtener el código y configurar", MUTED),
        ("git clone <repositorio> && cd bbti-erp", TEXT),
        ("cp .env.production.example .env.production", TEXT),
        ("", TEXT),
        ("# 3 · Levantar el sistema completo", MUTED),
        ("docker compose up -d", TEXT),
        ("", TEXT),
        ("# 4 · Cargar los datos iniciales (una vez)", MUTED),
        ("docker compose exec bbti-erp node prisma/seed.mjs", TEXT),
    ]
    y = 2.62
    for linea, color in comandos:
        if linea:
            texto(s, MARGIN + 0.32, y, ancho_izq - 0.6, 0.26, linea, size=10.5,
                  color=color, font=MONO)
        y += 0.275

    x_der = MARGIN + ancho_izq + 0.35
    ancho_der = W - MARGIN - x_der
    notas = [
        ("Sin intervención manual", "La base de datos se actualiza sola al arrancar y el certificado HTTPS se renueva automáticamente."),
        ("Respaldos diarios", "Copia automática de datos y archivos cada madrugada, con 30 días de historial y restauración probada."),
        ("Verificación continua", "Cada cambio del código dispara automáticamente todas las pruebas antes de aceptarse."),
    ]
    y = 2.35
    for tit, desc in notas:
        panel(s, x_der, y, ancho_der, 1.12)
        barra = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_der), Inches(y),
                                   Emu(28000), Inches(1.12))
        barra.fill.solid()
        barra.fill.fore_color.rgb = AMBER
        barra.line.fill.background()
        barra.shadow.inherit = False
        texto(s, x_der + 0.28, y + 0.18, ancho_der - 0.5, 0.3, tit, size=13, bold=True)
        texto(s, x_der + 0.28, y + 0.52, ancho_der - 0.5, 0.55, desc, size=10.5,
              color=MUTED, interlinea=1.3)
        y += 1.22
    pie(s, n)
    return s


def cierre(prs):
    s = lamina(prs)
    banda = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.14), prs.slide_height)
    banda.fill.solid()
    banda.fill.fore_color.rgb = AMBER
    banda.line.fill.background()
    banda.shadow.inherit = False

    texto(s, MARGIN, 2.85, 11, 1.4,
          "Sistema entregado.\nListo para producción.",
          size=38, bold=True, interlinea=1.1)
    regla(s, 4.65, w=3.2)
    texto(s, MARGIN, 4.95, 9.5, 0.9,
          "Queda a disposición para la puesta en marcha en cuanto se defina "
          "el servidor y el dominio.", size=15, color=MUTED, interlinea=1.4)
    texto(s, MARGIN, 6.3, 8, 0.3, "BBTI · SISTEMA DE GESTIÓN DE PROYECTOS",
          size=10, color=WAIT, font=MONO, space=1.6)
    return s


def main():
    salida = sys.argv[1] if len(sys.argv) > 1 else "BBTI-ERP-Estado.pptx"
    prs = nueva()

    portada(prs)
    resumen(prs, 2)
    alcance(prs, 3)
    migracion(prs, 4)
    lista_dos_columnas(
        prs, 5, "Posterior a la migración", "Mejoras incorporadas", None,
        [
            ("Pantallas que se actualizan solas.", "Lo que un usuario cambia aparece en la pantalla de los demás sin recargar."),
            ("Auditoría de producción.", "Cada sub-etapa guarda quién la confirmó, con fecha y hora de Lima."),
            ("Cambio de contraseña propio.", "Cada usuario la cambia desde la barra superior, sin pasar por el administrador."),
            ("Límite de archivos configurable.", "Ajustable sin tocar código, para cuando los planos pesen más."),
            ("Validación del acceso a datos.", "Herramienta que comprueba conexión y permisos sin exponer credenciales."),
            ("Verificación continua.", "Cada cambio del código dispara todas las pruebas automáticamente."),
        ], size_item=12)
    lista_dos_columnas(
        prs, 6, "Seguridad", "Protección aplicada y comprobada",
        "Diez medidas, todas verificadas con pruebas automáticas sobre el sistema real.",
        [
            ("Contraseñas cifradas", "y de 12 caracteres mínimo."),
            ("Bloqueo tras 5 intentos", "fallidos de acceso."),
            ("Sesiones revocables:", "cambiar la clave anula sesiones abiertas en otros equipos."),
            ("Baja inmediata:", "desactivar a una persona le corta el acceso en el acto."),
            ("Defensa contra suplantación", "del sistema dentro de sitios falsos (phishing)."),
            ("Verificación de origen", "en toda modificación de datos."),
            ("Validación estricta", "de todo dato que entra al sistema."),
            ("Permisos revalidados", "en el servidor en cada petición, no solo en pantalla."),
            ("Bitácora de seguridad:", "accesos, bloqueos y cambios de clave, con su IP."),
            ("Documentos privados,", "accesibles solo con enlaces firmados que caducan."),
        ], size_item=11)
    verificacion(prs, 7)
    estado(prs, 8)
    pendiente(prs, 9)
    despliegue(prs, 10)
    cierre(prs)

    prs.save(salida)
    print(f"Presentación generada: {salida} ({len(prs.slides.__iter__.__self__._sldIdLst)} láminas)")


if __name__ == "__main__":
    main()
